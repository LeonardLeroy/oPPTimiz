# oPPTimiz tests.
# Copyright (C) 2025 EDF. GNU General Public License v3 or later.

import io
import os
import zipfile
from xml.etree import ElementTree

import pytest
from PIL import Image

from opptimiz import optimize_bytes, optimize_file


def _big_jpeg(width: int = 4000, height: int = 3000) -> bytes:
    """A large, noisy JPEG that does not compress well, so recompression has work to do."""
    image = Image.frombytes("RGB", (width, height), os.urandom(width * height * 3))
    buffer = io.BytesIO()
    image.save(buffer, format="JPEG", quality=98)
    return buffer.getvalue()


@pytest.fixture
def real_pptx() -> bytes:
    pptx = pytest.importorskip("pptx")
    from pptx.util import Inches

    presentation = pptx.Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_picture(io.BytesIO(_big_jpeg()), Inches(0), Inches(0), Inches(10), Inches(7.5))

    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


@pytest.fixture
def real_docx() -> bytes:
    docx = pytest.importorskip("docx")

    document = docx.Document()
    document.add_picture(io.BytesIO(_big_jpeg()), width=docx.shared.Inches(6))

    buffer = io.BytesIO()
    document.save(buffer)
    return buffer.getvalue()


@pytest.fixture
def real_xlsx(tmp_path) -> bytes:
    openpyxl = pytest.importorskip("openpyxl")
    from openpyxl.drawing.image import Image as XlsxImage

    image_path = tmp_path / "photo.jpg"
    image_path.write_bytes(_big_jpeg())

    workbook = openpyxl.Workbook()
    workbook.active.add_image(XlsxImage(str(image_path)), "A1")

    buffer = io.BytesIO()
    workbook.save(buffer)
    return buffer.getvalue()


def _synthetic_ooxml(media: dict[str, bytes], rels_target: str) -> bytes:
    buffer = io.BytesIO()
    with zipfile.ZipFile(buffer, "w", zipfile.ZIP_DEFLATED) as archive:
        archive.writestr(
            "[Content_Types].xml",
            '<?xml version="1.0"?><Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">'
            '<Default Extension="png" ContentType="image/png"/>'
            '<Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>'
            '<Default Extension="xml" ContentType="application/xml"/></Types>',
        )
        archive.writestr("ppt/slides/slide1.xml", '<?xml version="1.0"?><sld><pic embed="rId1"/></sld>')
        archive.writestr(
            "ppt/slides/_rels/slide1.xml.rels",
            '<?xml version="1.0"?><Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
            f'<Relationship Id="rId1" Type="image" Target="../media/{rels_target}"/></Relationships>',
        )
        for name, data in media.items():
            archive.writestr(f"ppt/media/{name}", data)
    return buffer.getvalue()


def test_real_pptx_gets_smaller(real_pptx):
    optimized, result = optimize_bytes(real_pptx, level="maximal")
    assert result.new_size < result.initial_size
    assert result.gain > 0
    assert result.images_optimized >= 1


def test_real_pptx_still_opens(real_pptx):
    pptx = pytest.importorskip("pptx")
    optimized, _ = optimize_bytes(real_pptx, level="maximal")
    reopened = pptx.Presentation(io.BytesIO(optimized))
    assert len(reopened.slides) == 1


def test_output_is_a_valid_zip(real_pptx):
    optimized, _ = optimize_bytes(real_pptx, level="maximal")
    with zipfile.ZipFile(io.BytesIO(optimized)) as archive:
        assert archive.testzip() is None
        assert "[Content_Types].xml" in archive.namelist()


def test_real_docx_gets_smaller_and_still_opens(real_docx):
    docx = pytest.importorskip("docx")
    optimized, result = optimize_bytes(real_docx, level="maximal")
    assert result.new_size < result.initial_size
    assert result.images_optimized >= 1
    reopened = docx.Document(io.BytesIO(optimized))
    assert reopened.inline_shapes


def test_real_xlsx_gets_smaller_and_still_opens(real_xlsx):
    openpyxl = pytest.importorskip("openpyxl")
    optimized, result = optimize_bytes(real_xlsx, level="maximal")
    assert result.new_size < result.initial_size
    assert result.images_optimized >= 1
    openpyxl.load_workbook(io.BytesIO(optimized))


def test_unused_media_is_pruned():
    png = io.BytesIO()
    Image.new("RGB", (10, 10), "red").save(png, format="PNG")
    document = _synthetic_ooxml(
        {"image1.png": png.getvalue(), "orphan.png": png.getvalue()},
        rels_target="image1.png",
    )
    optimized, result = optimize_bytes(document, prune_unused=True)
    with zipfile.ZipFile(io.BytesIO(optimized)) as archive:
        names = archive.namelist()
    assert "ppt/media/image1.png" in names
    assert "ppt/media/orphan.png" not in names
    assert result.media_removed == 1


def test_no_prune_keeps_orphans():
    png = io.BytesIO()
    Image.new("RGB", (10, 10), "blue").save(png, format="PNG")
    document = _synthetic_ooxml(
        {"image1.png": png.getvalue(), "orphan.png": png.getvalue()},
        rels_target="image1.png",
    )
    optimized, result = optimize_bytes(document, prune_unused=False)
    with zipfile.ZipFile(io.BytesIO(optimized)) as archive:
        assert "ppt/media/orphan.png" in archive.namelist()
    assert result.media_removed == 0


def test_non_media_parts_are_preserved_bit_for_bit():
    png = io.BytesIO()
    Image.new("RGB", (10, 10), "green").save(png, format="PNG")
    document = _synthetic_ooxml({"image1.png": png.getvalue()}, rels_target="image1.png")

    with zipfile.ZipFile(io.BytesIO(document)) as archive:
        before = {n: archive.read(n) for n in archive.namelist() if "/media/" not in n}

    optimized, _ = optimize_bytes(document)
    with zipfile.ZipFile(io.BytesIO(optimized)) as archive:
        after = {n: archive.read(n) for n in archive.namelist() if "/media/" not in n}

    assert before == after
    for name, data in after.items():
        if name.endswith(".xml") or name.endswith(".rels"):
            ElementTree.fromstring(data)


def test_keep_original_writes_suffixed_file(tmp_path, real_pptx):
    source = tmp_path / "deck.pptx"
    source.write_bytes(real_pptx)

    destination, result = optimize_file(str(source), keep_original=True)

    assert destination == str(tmp_path / "deck_oPPTimiz.pptx")
    assert source.read_bytes() == real_pptx
    assert (tmp_path / "deck_oPPTimiz.pptx").exists()
    assert result.new_size < result.initial_size


def test_unsupported_extension_raises(tmp_path):
    bogus = tmp_path / "notes.txt"
    bogus.write_text("hello")
    with pytest.raises(ValueError):
        optimize_file(str(bogus))
